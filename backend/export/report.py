from __future__ import annotations
from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable,
)

_NAVY  = RGBColor(0x1A, 0x3A, 0x5C)
_TEAL  = RGBColor(0x08, 0x91, 0xB2)
_GRAY  = RGBColor(0x64, 0x74, 0x8B)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def build_pptx_bytes(eda_result: dict, narrative: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _add_title_slide(prs, blank, eda_result)
    _add_quality_slide(prs, blank, eda_result)
    _add_narrative_slide(prs, blank, narrative)
    _add_recommendations_slide(prs, blank, narrative)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _txt(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color is not None else _WHITE


def _add_title_slide(prs, layout, eda):
    s = prs.slides.add_slide(layout)
    _bg(s, 0x1A, 0x3A, 0x5C)
    summary = eda.get("summary", {})
    source  = summary.get("source_type", "Unknown").upper()
    rows    = summary.get("total_rows", "N/A")
    today   = date.today().strftime("%d %B %Y")
    _txt(s, "SmartDQC", 0.5, 0.4, 8, 0.5, size=14, bold=True, color=_TEAL)
    _txt(s, "Data Quality Report", 0.5, 0.9, 10, 1.0, size=36, bold=True)
    _txt(s, f"Source: {source}  |  Records: {rows}  |  {today}",
         0.5, 2.0, 11, 0.5, size=14, color=RGBColor(0xA8, 0xC8, 0xD8))


def _add_quality_slide(prs, layout, eda):
    s = prs.slides.add_slide(layout)
    _bg(s, 0xF0, 0xF7, 0xFA)
    _txt(s, "Data Quality Overview", 0.5, 0.3, 10, 0.6,
         size=24, bold=True, color=_NAVY)
    quality    = eda.get("quality", {})
    indicators = eda.get("indicators", {})
    outliers   = eda.get("outliers", {})
    lines = [
        f"Overall Quality Score : {quality.get('overall_score', 'N/A')}",
        f"Completeness          : {quality.get('overall_completeness', 'N/A')}%",
        f"Missing Data Rate     : {quality.get('missing_rate', 'N/A')}",
        f"Outliers Flagged      : {outliers.get('total_flagged', 'N/A')}",
    ]
    for flag in ["stunting_rate", "wasting_rate", "underweight_rate", "overweight_rate"]:
        if flag in indicators:
            label = flag.replace("_rate", "").capitalize()
            val   = indicators[flag]
            if isinstance(val, float):
                val = round(val * 100, 1)
            lines.append(f"{label:22s}: {val}%")
    _txt(s, "\n".join(lines), 0.5, 1.1, 12, 5.5, size=14, color=_NAVY)


def _add_narrative_slide(prs, layout, narrative):
    s = prs.slides.add_slide(layout)
    _bg(s, 0x1A, 0x3A, 0x5C)
    _txt(s, "AI Analysis Summary", 0.5, 0.3, 10, 0.6, size=24, bold=True)
    exec_sum = narrative.get("executive_summary", {})
    _txt(s, "Bahasa Malaysia", 0.5, 1.1, 5, 0.4, size=11, bold=True, color=_TEAL)
    _txt(s, exec_sum.get("bm", "–"), 0.5, 1.55, 5.8, 4.5, size=12)
    _txt(s, "English", 6.9, 1.1, 5, 0.4, size=11, bold=True, color=_TEAL)
    _txt(s, exec_sum.get("en", "–"), 6.9, 1.55, 5.8, 4.5, size=12)


def _add_recommendations_slide(prs, layout, narrative):
    s = prs.slides.add_slide(layout)
    _bg(s, 0xF0, 0xF7, 0xFA)
    _txt(s, "Recommendations", 0.5, 0.3, 10, 0.6,
         size=24, bold=True, color=_NAVY)
    for i, rec in enumerate(narrative.get("recommendations", [])[:3]):
        y = 1.1 + i * 1.9
        priority = rec.get("priority", "").upper()
        _txt(s, f"[{priority}] {rec.get('action', '')}", 0.5, y, 12, 0.45,
             size=13, bold=True, color=_NAVY)
        _txt(s, rec.get("en", ""), 0.5, y + 0.45, 12, 1.3, size=11, color=_GRAY)


# ── PDF ───────────────────────────────────────────────────────────────────────

def build_pdf_bytes(eda_result: dict, narrative: dict) -> bytes:
    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    h1  = ParagraphStyle("H1",  parent=styles["Heading1"],  fontSize=20,
                         textColor=rl_colors.HexColor("#1A3A5C"))
    h2  = ParagraphStyle("H2",  parent=styles["Heading2"],  fontSize=14,
                         textColor=rl_colors.HexColor("#0891B2"))
    bod = ParagraphStyle("Body", parent=styles["Normal"],   fontSize=11, leading=16)
    sm  = ParagraphStyle("Sm",   parent=styles["Normal"],   fontSize=9,
                         textColor=rl_colors.HexColor("#64748B"))

    story = []
    today   = date.today().strftime("%d %B %Y")
    summary = eda_result.get("summary", {})

    story.append(Paragraph("SmartDQC — Data Quality Report", h1))
    story.append(Paragraph(
        f"Source: {summary.get('source_type','N/A').upper()}  |  "
        f"Records: {summary.get('total_rows','N/A')}  |  {today}", sm))
    story.append(HRFlowable(width="100%", thickness=2,
                            color=rl_colors.HexColor("#0891B2")))
    story.append(Spacer(1, 0.4*cm))

    story.append(Paragraph("Data Quality Overview", h2))
    quality    = eda_result.get("quality", {})
    indicators = eda_result.get("indicators", {})
    outliers   = eda_result.get("outliers", {})
    q_rows = [
        ["Metric", "Value"],
        ["Overall Quality Score",  str(quality.get("overall_score", "N/A"))],
        ["Completeness",           f"{quality.get('overall_completeness','N/A')}%"],
        ["Missing Data Rate",      str(quality.get("missing_rate", "N/A"))],
        ["Outliers Flagged",       str(outliers.get("total_flagged", "N/A"))],
    ]
    for flag in ["stunting_rate", "wasting_rate", "underweight_rate", "overweight_rate"]:
        if flag in indicators:
            val = indicators[flag]
            if isinstance(val, float):
                val = round(val * 100, 1)
            q_rows.append([flag.replace("_rate", "").capitalize(), f"{val}%"])

    tbl = Table(q_rows, colWidths=[10*cm, 6*cm])
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0), rl_colors.HexColor("#1A3A5C")),
        ("TEXTCOLOR",     (0, 0), (-1, 0), rl_colors.white),
        ("FONTSIZE",      (0, 0), (-1, -1), 10),
        ("GRID",          (0, 0), (-1, -1), 0.5, rl_colors.HexColor("#E2EEF4")),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1),
         [rl_colors.white, rl_colors.HexColor("#F0F7FA")]),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING",    (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(tbl)
    story.append(Spacer(1, 0.6*cm))

    exec_sum = narrative.get("executive_summary", {})
    if exec_sum:
        story.append(Paragraph("AI Analysis Summary", h2))
        story.append(Paragraph(
            f"<b>Bahasa Malaysia</b><br/>{exec_sum.get('bm','')}", bod))
        story.append(Spacer(1, 0.3*cm))
        story.append(Paragraph(
            f"<b>English</b><br/>{exec_sum.get('en','')}", bod))
        story.append(Spacer(1, 0.6*cm))

    recs = narrative.get("recommendations", [])
    if recs:
        story.append(Paragraph("Recommendations", h2))
        for rec in recs[:5]:
            story.append(Paragraph(
                f"<b>[{rec.get('priority','').upper()}] {rec.get('action','')}</b>", bod))
            story.append(Paragraph(rec.get("en", ""), bod))
            story.append(Spacer(1, 0.2*cm))

    doc.build(story)
    return buf.getvalue()
